# -*- coding: utf-8 -*-
"""
bongbong7 팀 - 제8회 공군 해커톤 AI경진대회 본선 발표자료 (전체 12슬라이드)
디자인: 화이트 배경 + 틸→네이비 그라데이션 램프, 프로세스/타임라인 인포그래픽 스타일
평가기준(발표평가 20점: 문제해결접근성5·AI모델설계5·임무수행전략5·발표및질의응답5) 대응 구성
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

FONT = "맑은 고딕"

# ---------------- 색상 팔레트 (라이트 · 틸/네이비) ----------------
BG        = RGBColor(0xF6, 0xF9, 0xFA)
CARD_BG   = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BD   = RGBColor(0xE2, 0xE9, 0xEC)
ROW_ALT   = RGBColor(0xF3, 0xF7, 0xF8)
LINE      = RGBColor(0xDC, 0xE4, 0xE8)
TEAL      = RGBColor(0x25, 0x8E, 0xA6)
TEAL_DK   = RGBColor(0x14, 0x62, 0x77)
NAVY      = RGBColor(0x17, 0x33, 0x4C)
NAVY_DEEP = RGBColor(0x11, 0x26, 0x38)
TEAL_TINT = RGBColor(0xEA, 0xF6, 0xF7)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY = RGBColor(0x5B, 0x6B, 0x79)
TEXT_MUTE = RGBColor(0x93, 0xA1, 0xAC)
ORANGE_BG = RGBColor(0xFC, 0xEE, 0xDA)
ORANGE_TX = RGBColor(0xB0, 0x63, 0x0B)
TEAL_BG2  = RGBColor(0xDD, 0xF1, 0xF3)
TEAL_TX2  = RGBColor(0x0E, 0x5C, 0x6E)


def lerp_color(c1, c2, t):
    r = int(c1[0] + (c2[0] - c1[0]) * t)
    g = int(c1[1] + (c2[1] - c1[1]) * t)
    b = int(c1[2] + (c2[2] - c1[2]) * t)
    return RGBColor(r, g, b)


RAMP_START = (0x5B, 0xC6, 0xCF)
RAMP_END   = (0x17, 0x33, 0x4C)


def ramp(i, n):
    t = i / max(1, n - 1)
    return lerp_color(RAMP_START, RAMP_END, t)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

MX = Inches(0.55)
CW = Inches(12.23)


# ---------------- 헬퍼 함수 ----------------
def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75,
             shape_type=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None and shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def add_icon(slide, shape_type, x, y, w, h, fill=WHITE, line=None):
    sp = slide.shapes.add_shape(shape_type, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def add_ring(slide, x, y, d, color, w_pt=2.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    sp.fill.background()
    sp.line.color.rgb = color
    sp.line.width = Pt(w_pt)
    sp.shadow.inherit = False
    return sp


def add_text(slide, x, y, w, h, text, size=12, color=NAVY, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
             line_spacing=1.0, font=FONT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line_text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line_text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color
        rPr = r._r.get_or_add_rPr()
        ea = rPr.makeelement(qn('a:ea'), {'typeface': font})
        rPr.append(ea)
    return tb


def add_pill(slide, x, y, w, h, text, fill, text_color=WHITE, size=8.5, bold=True):
    sp = add_rect(slide, x, y, w, h, fill=fill,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = sp.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = FONT
    r.font.color.rgb = text_color
    return sp


def header_footer(slide, page_no, foot_note, eyebrow):
    add_text(slide, MX, Inches(0.30), Inches(8), Inches(0.3),
              "bongbong7  ·  제8회 공군 해커톤 AI경진대회 본선 발표", size=10, color=TEXT_MUTE)
    add_text(slide, Inches(11.7), Inches(0.24), Inches(1.05), Inches(0.4),
              page_no, size=17, color=TEAL, bold=True, align=PP_ALIGN.RIGHT)
    add_rect(slide, MX, Inches(0.64), CW, Pt(1.2), fill=LINE)

    add_rect(slide, MX, Inches(7.08), CW, Pt(1.2), fill=LINE)
    add_text(slide, MX, Inches(7.16), Inches(8), Inches(0.3), foot_note, size=8.5, color=TEXT_MUTE)
    add_icon(slide, MSO_SHAPE.DIAMOND, Inches(11.55), Inches(7.17), Inches(0.11), Inches(0.11), fill=TEAL)
    add_text(slide, Inches(11.74), Inches(7.14), Inches(1.0), Inches(0.3),
              "bongbong7", size=9.5, color=NAVY, bold=True)

    add_pill(slide, MX, Inches(0.86), Inches(2.15), Inches(0.3), eyebrow,
             fill=NAVY, text_color=WHITE, size=9)


def title_block(slide, title, subtitle):
    add_text(slide, MX, Inches(1.26), CW, Inches(0.48), title, size=23, color=NAVY, bold=True)
    add_text(slide, MX, Inches(1.76), CW, Inches(0.32), subtitle, size=12, color=TEXT_GRAY)
    add_rect(slide, MX, Inches(2.16), CW, Pt(1.2), fill=LINE)


def chevron(slide, cx, cy, size=0.19, color=TEAL):
    w = Inches(size * 1.35)
    h = Inches(size * 1.7)
    add_icon(slide, MSO_SHAPE.CHEVRON, cx - w / 2, cy - h / 2, w, h, fill=color)


def numbered_bullet(slide, x, y, w, num_color, text, size=8.6, text_color=NAVY, d=0.22):
    add_icon(slide, MSO_SHAPE.OVAL, x, y, Inches(d), Inches(d), fill=num_color)
    add_text(slide, x, y, Inches(d), Inches(d), "", size=1)  # placeholder no-op
    add_text(slide, x + Inches(d) + Inches(0.12), y - Inches(0.02), w - Inches(d) - Inches(0.12), Inches(0.3),
              text, size=size, color=text_color, anchor=MSO_ANCHOR.MIDDLE)


# =====================================================================
# SLIDE 01 : 표지
# =====================================================================
s1 = prs.slides.add_slide(blank)
set_bg(s1, BG)

# 장식용 링(추상 원형 액센트)
add_ring(s1, Inches(10.6), Inches(-1.1), Inches(4.2), ramp(1, 5), 2.2)
add_ring(s1, Inches(11.6), Inches(0.4), Inches(2.0), ramp(3, 5), 1.6)
add_ring(s1, Inches(-1.0), Inches(5.3), Inches(3.4), ramp(4, 5), 1.8)
sp = add_icon(s1, MSO_SHAPE.OVAL, Inches(11.95), Inches(1.55), Inches(0.22), Inches(0.22), fill=ramp(2, 5))

add_pill(s1, MX, Inches(1.05), Inches(3.3), Inches(0.34),
         "제8회 공군 해커톤 AI경진대회 · 본선 발표", fill=NAVY, text_color=WHITE, size=9.5)

add_text(s1, MX, Inches(2.55), Inches(10.6), Inches(0.75),
          "AI 기반 정찰·전투피해평가(BDA)", size=34, color=NAVY, bold=True)
add_text(s1, MX, Inches(3.28), Inches(10.6), Inches(0.75),
          "자동화 파이프라인", size=34, color=TEAL_DK, bold=True)

add_rect(s1, MX, Inches(4.18), Inches(0.5), Pt(4), fill=TEAL)
add_text(s1, MX, Inches(4.35), Inches(9.5), Inches(0.4),
          "영상 한 장으로 8종 결과 JSON까지 — 좌표계 통일 기반 6단계 완전자동 파이프라인",
          size=13, color=TEXT_GRAY)

# 팀 정보 뱃지 행
info_y = Inches(5.15)
add_pill(s1, MX, info_y, Inches(1.9), Inches(0.4), "Team bongbong7", fill=TEAL, text_color=WHITE, size=10.5)
add_text(s1, MX + Inches(2.05), info_y + Inches(0.02), Inches(3.0), Inches(0.36),
          "2026. 07. 15  본선 발표", size=11, color=TEXT_GRAY, anchor=MSO_ANCHOR.MIDDLE)

# 하단 소형 램프 도트 (브랜드 플로리시)
dot_y = Inches(6.55)
labels = ["영상입력", "좌표정합", "AI탐지·분류", "결과통합", "LLM보고·전송"]
dot_x0 = MX
gap = Inches(2.3)
for i, lb in enumerate(labels):
    dx = dot_x0 + i * gap
    add_icon(s1, MSO_SHAPE.OVAL, dx, dot_y, Inches(0.16), Inches(0.16), fill=ramp(i, 5))
    add_text(s1, dx + Inches(0.24), dot_y - Inches(0.02), Inches(2.0), Inches(0.22),
              lb, size=9, color=TEXT_GRAY)
    if i < 4:
        add_rect(s1, dx + Inches(1.05), dot_y + Inches(0.07), Inches(1.05), Pt(1), fill=LINE)


# =====================================================================
# SLIDE 02 : 파이프라인 개요 (프로세스 인포그래픽 · 2행 4열 그리드)
# =====================================================================
s2 = prs.slides.add_slide(blank)
set_bg(s2, BG)
header_footer(s2, "02", "config/field_config.py · src/pipeline.py 기준", "PIPELINE OVERVIEW")
title_block(
    s2,
    "드론 영상 한 장을 8종 결과 JSON으로 — 6단계 자동 파이프라인",
    "영상 입력 → 좌표 정합 → AI 탐지·분류 → 결과통합 → LLM 보고 → 대시보드 전송까지, 사람 개입 없이 순차 자동 실행"
)

steps = [
    ("STEP 01", "영상 입력", "watchdog 자동 감지 · frame 주기 추출", MSO_SHAPE.ISOSCELES_TRIANGLE, "video_watcher.py"),
    ("STEP 02", "탑뷰 워핑", "ArUco 검출 → 호모그래피 변환", MSO_SHAPE.DIAMOND, "calibration.py"),
    ("STEP 03", "활주로·유도로 탐지", "20개 zone 타일 · YOLO11n 배치추론", MSO_SHAPE.DONUT, "tiling.py / detection.py"),
    ("STEP 04", "시설물 상태 분류", "고정 6개소 ROI · 정상·파손·화재", MSO_SHAPE.PENTAGON, "detection.py"),
    ("STEP 05", "결과 통합", "실좌표 중복제거 · 8종 JSON 매핑", MSO_SHAPE.HEXAGON, "geo_dedup.py / schemas.py"),
    ("STEP 06", "상황보고서", "LLM 프롬프트 · 50~100자 생성", MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT, "report_generator.py"),
    ("STEP 07", "대시보드 전송", "자동 QA 검증 · 순차 전송", MSO_SHAPE.UP_ARROW, "validator.py / transmitter.py"),
]

GRID_Y0 = Inches(2.55)
ROW_H = Inches(2.05)
ROW_GAP = Inches(0.42)
COL_GAP = Inches(0.24)
NCOL = 4
col_w = int((CW - COL_GAP * (NCOL - 1)) / NCOL)

row1_y = GRID_Y0
row2_y = GRID_Y0 + ROW_H + ROW_GAP


def col_x(i):
    return MX + i * (col_w + COL_GAP)


def process_card(slide, x, y, w, h, badge_txt, title, desc, shape, fname, color):
    add_rect(slide, x, y, w, h, fill=CARD_BG, line=CARD_BD, line_w=0.75,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.055)
    pad = Inches(0.16)
    dia = Inches(0.56)
    add_icon(slide, MSO_SHAPE.OVAL, x + pad, y + pad, dia, dia, fill=color)
    add_icon(slide, shape, x + pad + Inches(0.14), y + pad + Inches(0.14),
              dia - Inches(0.28), dia - Inches(0.28), fill=WHITE)
    badge_w = Inches(0.86)
    add_pill(slide, x + w - badge_w - pad, y + pad + Inches(0.03), badge_w, Inches(0.24),
              badge_txt, fill=TEAL_TINT, text_color=TEAL_DK, size=8.3, bold=True)
    add_text(slide, x + pad, y + pad + dia + Inches(0.10), w - pad * 2, Inches(0.55),
              title, size=13, color=NAVY, bold=True, line_spacing=1.05)
    add_text(slide, x + pad, y + pad + dia + Inches(0.62), w - pad * 2, Inches(0.55),
              desc, size=9.6, color=TEXT_GRAY, line_spacing=1.15)
    add_rect(slide, x + pad, y + h - Inches(0.34), w - pad * 2, Inches(0.24),
              fill=TEAL_TINT, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    add_text(slide, x + pad, y + h - Inches(0.32), w - pad * 2, Inches(0.2),
              fname, size=7.6, color=TEAL_DK, align=PP_ALIGN.CENTER)


for i in range(4):
    badge, title, desc, shape, fname = steps[i]
    x = col_x(i)
    process_card(s2, x, row1_y, col_w, ROW_H, badge, title, desc, shape, fname, ramp(i, 7))
    if i < 3:
        chevron(s2, x + col_w + COL_GAP / 2, row1_y + ROW_H / 2, size=0.19, color=TEAL)

for i in range(3):
    idx = 4 + i
    badge, title, desc, shape, fname = steps[idx]
    x = col_x(i)
    process_card(s2, x, row2_y, col_w, ROW_H, badge, title, desc, shape, fname, ramp(idx, 7))
    chevron(s2, x + col_w + COL_GAP / 2, row2_y + ROW_H / 2, size=0.19, color=TEAL)

cx = col_x(3)
add_rect(s2, cx, row2_y, col_w, ROW_H, fill=NAVY_DEEP,
          shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.055)
pad = Inches(0.16)
dia = Inches(0.42)
add_icon(s2, MSO_SHAPE.OVAL, cx + pad, row2_y + pad, dia, dia, fill=TEAL)
add_icon(s2, MSO_SHAPE.STAR_5_POINT, cx + pad + Inches(0.08), row2_y + pad + Inches(0.08),
          dia - Inches(0.16), dia - Inches(0.16), fill=WHITE)
add_text(s2, cx + pad, row2_y + pad + dia + Inches(0.05), col_w - pad * 2, Inches(0.26),
          "왜 이렇게 설계했나", size=11.5, color=WHITE, bold=True)
bullets = [
    "20개 zone · 배치추론 1회 처리",
    "시설물 6슬롯 강제 매핑",
    "로컬 LLM + 오프라인 폴백",
    "전송 전 자동 QA 검증",
]
by = row2_y + pad + dia + Inches(0.36)
for b in bullets:
    add_icon(s2, MSO_SHAPE.OVAL, cx + pad, by + Inches(0.045), Inches(0.045), Inches(0.045), fill=TEAL)
    add_text(s2, cx + pad + Inches(0.13), by, col_w - pad * 2 - Inches(0.13), Inches(0.2),
              b, size=8.3, color=WHITE, line_spacing=1.0)
    by += Inches(0.225)


# =====================================================================
# SLIDE 03 : 개발 타임라인 (git 커밋 기반 · 지그재그 타임라인)
# =====================================================================
s3 = prs.slides.add_slide(blank)
set_bg(s3, BG)
header_footer(s3, "03", "git log 2026-07-14 08:12 → 07-15 11:06 기준", "DEVELOPMENT LOG")
title_block(
    s3,
    "27시간, 73번의 커밋이 만든 파이프라인",
    "2026.07.14 08:12 ~ 07.15 11:06  ·  커밋 73개  ·  PR 병합 24건  ·  현장에서 발견한 문제를 그 자리에서 고친 기록"
)

phases = [
    ("Day 1 · 오전", "08:12–11:15", "6단계 BDA 아키텍처 확정",
     "8종 JSON 규격 미비", "좌표계 통일로 재설계", "up"),
    ("Day 1 · 오후", "12:07–15:48", "8종 JSON 스키마 재작성",
     "ArUco 마커ID 실측 불일치", "근접크롭으로 좌표 재정의", "down"),
    ("Day 1 · 저녁", "16:52–18:32", "대시보드 연동 현장 검증",
     "전송 스펙, 실서버와 상이", "전송 모듈 즉시 재작성", "up"),
    ("Day 1 · 심야", "20:03–21:41", "LLM 상황보고 100자 규정화",
     "UXO 탐지 개수 초과 집계", "구간별 상한 강제 적용", "down"),
    ("Day 2 · 오전", "07:57–11:06", "크레이터 크기 재판정",
     "크기 오분류·중복 위험", "라벨 보정 + 앙상블 검증", "up"),
]

LINE_Y = Inches(4.68)
add_rect(s3, MX + Inches(0.3), LINE_Y, CW - Inches(0.6), Pt(2), fill=TEAL)

n = len(phases)
slot_w = CW / n
NODE_D = Inches(0.44)
CARD_W = Inches(2.32)
CARD_H = Inches(1.78)
STEM_LEN = Inches(0.24)


def tag_row(slide, x, y, w, label, fill, text_color, desc, desc_color):
    pill_w = Inches(0.46)
    add_pill(slide, x, y, pill_w, Inches(0.2), label, fill=fill, text_color=text_color, size=7.3)
    add_text(slide, x + pill_w + Inches(0.08), y, w - pill_w - Inches(0.08), Inches(0.2),
              desc, size=7.6, color=desc_color, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)


for i, (phase, trange, work, problem, fix, direction) in enumerate(phases):
    cx = MX + slot_w * i + slot_w / 2
    color = ramp(i, n)

    add_icon(s3, MSO_SHAPE.OVAL, cx - NODE_D / 2, LINE_Y - NODE_D / 2 + Pt(1), NODE_D, NODE_D, fill=color)
    add_text(s3, cx - NODE_D / 2, LINE_Y - NODE_D / 2 + Pt(1), NODE_D, NODE_D,
              str(i + 1).zfill(2), size=12, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    card_x = cx - CARD_W / 2
    if direction == "up":
        stem_y0 = LINE_Y - NODE_D / 2 - STEM_LEN
        stem_y1 = LINE_Y - NODE_D / 2
        card_y = stem_y0 - CARD_H
    else:
        stem_y0 = LINE_Y + NODE_D / 2
        stem_y1 = LINE_Y + NODE_D / 2 + STEM_LEN
        card_y = stem_y1

    add_rect(s3, cx - Pt(1), stem_y0, Pt(2), stem_y1 - stem_y0, fill=color)

    add_rect(s3, card_x, card_y, CARD_W, CARD_H, fill=CARD_BG, line=CARD_BD, line_w=0.75,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.055)

    pad = Inches(0.14)
    inner_w = CARD_W - pad * 2
    add_text(s3, card_x + pad, card_y + Inches(0.09), inner_w, Inches(0.22),
              phase, size=11, color=NAVY, bold=True)
    add_text(s3, card_x + pad, card_y + Inches(0.32), inner_w, Inches(0.18),
              trange, size=8.2, color=TEAL_DK, bold=True)
    add_rect(s3, card_x + pad, card_y + Inches(0.54), inner_w, Pt(1), fill=LINE)

    add_text(s3, card_x + pad, card_y + Inches(0.60), inner_w, Inches(0.15),
              "핵심 작업", size=7.2, color=TEXT_MUTE, bold=True)
    add_text(s3, card_x + pad, card_y + Inches(0.76), inner_w, Inches(0.30),
              work, size=8.0, color=NAVY, line_spacing=1.08)

    tag_row(s3, card_x + pad, card_y + Inches(1.12), inner_w, "문제",
            ORANGE_BG, ORANGE_TX, problem, TEXT_GRAY)
    tag_row(s3, card_x + pad, card_y + Inches(1.38), inner_w, "해결",
            TEAL_BG2, TEAL_TX2, fix, NAVY)


# =====================================================================
# SLIDE 04 : 과제 이해 & 문제해결 접근 전략  (발표평가 "문제해결 접근성" 대응)
# =====================================================================
s4 = prs.slides.add_slide(blank)
set_bg(s4, BG)
header_footer(s4, "04", "대회 임무평가 80점 항목(사전교육자료 p.38) 기준", "PROBLEM APPROACH")
title_block(
    s4,
    "임무평가 80점, 5개 항목에 정면으로 대응합니다",
    "핵심 접근: ArUco로 좌표계를 통일해 픽셀좌표를 실좌표에 정렬 — 이후 모든 탐지·분류는 단순 슬라이싱으로 처리"
)

t_y = Inches(2.48)
cols = [
    ("임무 항목 (배점)", Inches(3.55)),
    ("난이도", Inches(0.95)),
    ("우리 팀의 해결 전략", Inches(5.35)),
    ("확정 로직", Inches(2.06)),
]
row_h = Inches(0.60)
header_h = Inches(0.46)

# 헤더 행
hx = MX
add_rect(s4, MX, t_y, CW, header_h, fill=NAVY)
for label, w in cols:
    add_text(s4, hx + Inches(0.14), t_y, w - Inches(0.14), header_h,
              label, size=10, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    hx += w

rows = [
    ("활주로·유도로 폭파구 탐지", "(25점)", "중",
     "20개 zone 배치추론 → 위치·크기·개수·가용길이 산출", "runway_analysis.py"),
    ("시설물 피해 확인", "(18점)", "상",
     "고정 6개 ROI 분류(정상/파손/화재) · 6슬롯 강제 매핑", "facility_analysis.py"),
    ("위험 장애물 탐지", "(20점)", "중",
     "미사일·포탄·자탄 탐지 + 구간별 개수 상한(cap) 적용", "uxo_analysis.py"),
    ("LLM 기반 상황보고", "(7점)", "상",
     "로컬 LLM + 오프라인 템플릿 이중화, 100자 제약 자동 준수", "report_generator.py"),
    ("임무 수행 시간", "(10점)", "중",
     "zone 타일 배치추론 1회 호출로 3분(180초) 제한 대응", "pipeline.py"),
]

ry = t_y + header_h
diff_color = {"중": (TEAL_BG2, TEAL_TX2), "상": (ORANGE_BG, ORANGE_TX)}
for i, (name, pt, diff, strat, fname) in enumerate(rows):
    bgc = CARD_BG if i % 2 == 0 else ROW_ALT
    add_rect(s4, MX, ry, CW, row_h, fill=bgc)
    rx = MX
    add_text(s4, rx + Inches(0.14), ry + Inches(0.08), cols[0][1] - Inches(0.14), Inches(0.24),
              name, size=9.9, color=NAVY, bold=True, line_spacing=1.0)
    add_text(s4, rx + Inches(0.14), ry + Inches(0.34), cols[0][1] - Inches(0.14), Inches(0.2),
              pt, size=8.2, color=TEAL_DK, bold=True)
    rx += cols[0][1]
    dfill, dtext = diff_color[diff]
    add_pill(s4, rx + Inches(0.10), ry + row_h / 2 - Inches(0.1), Inches(0.5), Inches(0.2),
             diff, fill=dfill, text_color=dtext, size=7.9)
    rx += cols[1][1]
    add_text(s4, rx + Inches(0.14), ry, cols[2][1] - Inches(0.24), row_h,
              strat, size=9.6, color=TEXT_GRAY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    rx += cols[2][1]
    add_rect(s4, rx + Inches(0.10), ry + row_h / 2 - Inches(0.12), cols[3][1] - Inches(0.24), Inches(0.24),
              fill=TEAL_TINT, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    add_text(s4, rx + Inches(0.10), ry + row_h / 2 - Inches(0.10), cols[3][1] - Inches(0.24), Inches(0.2),
              fname, size=7.9, color=TEAL_DK, align=PP_ALIGN.CENTER)
    ry += row_h

add_rect(s4, MX, t_y, CW, header_h + row_h * len(rows), fill=None, line=CARD_BD, line_w=1.0)


# =====================================================================
# SLIDE 05 : AI 모델 설계 & 차별화 전략  (발표평가 "AI 모델 설계" 대응)
# =====================================================================
s5 = prs.slides.add_slide(blank)
set_bg(s5, BG)
header_footer(s5, "05", "src/detection.py · calibration.py · geo_dedup.py 기준", "AI MODEL DESIGN")
title_block(
    s5,
    "표준 알고리즘 위에, 대회 규정을 겨냥한 7가지 설계 선택",
    "YOLO11n·ArUco 등 검증된 기술 위에 3분 제한·오프라인·8종 JSON 스펙에 맞춘 엔지니어링을 더했습니다"
)

# 좌측: 기술 스택 표
lt_x = MX
lt_w = Inches(5.35)
lt_y = Inches(2.55)
stack_rows = [
    ("폭파구 · 불발탄 탐지", "YOLO11n 객체탐지 + ArUco 호모그래피"),
    ("시설물 상태 분류", "YOLO11n-cls 이미지 분류"),
    ("좌표 정합", "OpenCV Homography / Perspective Transform"),
    ("상황보고 생성", "로컬 LLM(Ollama qwen2.5) + Prompt Engineering"),
]
add_text(s5, lt_x, lt_y, lt_w, Inches(0.3), "임무별 기술 스택", size=13, color=NAVY, bold=True)
sy = lt_y + Inches(0.42)
srow_h = Inches(0.62)
for i, (task, tech) in enumerate(stack_rows):
    bgc = CARD_BG if i % 2 == 0 else ROW_ALT
    add_rect(s5, lt_x, sy, lt_w, srow_h, fill=bgc, line=CARD_BD, line_w=0.6)
    add_icon(s5, MSO_SHAPE.OVAL, lt_x + Inches(0.14), sy + srow_h / 2 - Inches(0.06), Inches(0.12), Inches(0.12),
              fill=ramp(i, 4))
    add_text(s5, lt_x + Inches(0.38), sy + Inches(0.08), lt_w - Inches(0.5), Inches(0.24),
              task, size=9.9, color=NAVY, bold=True)
    add_text(s5, lt_x + Inches(0.38), sy + Inches(0.33), lt_w - Inches(0.5), Inches(0.26),
              tech, size=8.8, color=TEXT_GRAY)
    sy += srow_h

# 우측: 설계 철학 7가지
rt_x = Inches(6.30)
rt_w = Inches(6.48)
add_text(s5, rt_x, lt_y, rt_w, Inches(0.3), "설계 철학 7가지", size=13, color=NAVY, bold=True)
philosophy = [
    "좌표계 통일 + 탑뷰 워핑 — 역투영 없이 단순 슬라이싱으로 crop",
    "경계 없는 그리드 분할 — zone이 겹치지 않아 overlap 처리 불필요",
    "배치추론 1회 호출 — zone 타일·ROI를 한 번의 predict(list)로 처리",
    "실좌표 기반 중복 제거 — 같은 지점 다중 촬영에도 개수 안 부풀림",
    "시설물 6슬롯 강제 매핑 — 탐지 실패해도 unconfirmed로 항상 6개 보고",
    "로컬 LLM + 오프라인 폴백 — Ollama 실패해도 템플릿 자동 전환",
    "고전CV ↔ YOLO11n 백엔드 전환 — 현장 리스크 대비 안전장치",
]
py = lt_y + Inches(0.46)
for i, text in enumerate(philosophy):
    add_icon(s5, MSO_SHAPE.OVAL, rt_x, py, Inches(0.24), Inches(0.24), fill=ramp(i, 7))
    add_text(s5, rt_x, py, Inches(0.24), Inches(0.24), str(i + 1), size=9.5, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s5, rt_x + Inches(0.36), py - Inches(0.01), rt_w - Inches(0.36), Inches(0.34),
              text, size=9.4, color=NAVY, line_spacing=1.05, anchor=MSO_ANCHOR.MIDDLE)
    py += Inches(0.455)


# =====================================================================
# SLIDE 06 : 데이터 수집 전략
# =====================================================================
s6 = prs.slides.add_slide(blank)
set_bg(s6, BG)
header_footer(s6, "06", "bongbong7/배치.txt (팀 촬영 계획) 기준", "DATA COLLECTION")
title_block(
    s6,
    "촬영 계획 하나로 라벨 편향을 줄이는 데이터 수집 전략",
    "변동 요소(크기·배경·각도·상태)를 표로 미리 정리해, 빠짐없이 체계적으로 촬영합니다"
)

dc_y = Inches(2.55)
dc_gap = Inches(0.4)
dc_col_w = (CW - dc_gap) / 2
dc_lx = MX
dc_rx = MX + dc_col_w + dc_gap

add_text(s6, dc_lx, dc_y, dc_col_w, Inches(0.3), "폭파구 · 불발탄 촬영 전략", size=13, color=NAVY, bold=True)
crater_items = [
    ("정상 개별 근접 촬영", "폭파구·불발탄 각각 단독 클로즈업 먼저 확보(추후 증강용)"),
    ("크기 × 배경 조합 체계화", "대/중/소 × 배경(잔디/도로) 변동요소를 표로 정리"),
    ("각도 다양화 + 그림자 점검", "다양한 각도로 촬영하되 그림자 유무를 매번 확인"),
    ("배경 세부 유형 반영", "유도로 ㄱ·T·一자, 활주로 끝 흰색 마킹 2종까지 커버"),
    ("촬영 우선순위 설정", "폭파구+불발탄 동시 구도 우선, 서 있는 미사일 우선"),
]
iy = dc_y + Inches(0.44)
dc_step = Inches(0.62)
for i, (t, d) in enumerate(crater_items):
    add_icon(s6, MSO_SHAPE.OVAL, dc_lx, iy, Inches(0.26), Inches(0.26), fill=ramp(i, 5))
    add_text(s6, dc_lx, iy, Inches(0.26), Inches(0.26), str(i + 1), size=9.5, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s6, dc_lx + Inches(0.38), iy - Inches(0.02), dc_col_w - Inches(0.38), Inches(0.22),
              t, size=10.3, color=NAVY, bold=True)
    add_text(s6, dc_lx + Inches(0.38), iy + Inches(0.21), dc_col_w - Inches(0.38), Inches(0.22),
              d, size=8.8, color=TEXT_GRAY)
    iy += dc_step

add_text(s6, dc_rx, dc_y, dc_col_w, Inches(0.3), "시설물 촬영 전략", size=13, color=NAVY, bold=True)
facility_items = [
    ("정상 상태 전체 촬영 우선", "시설물 6종의 정상 상태 사진을 가장 먼저 확보"),
    ("45도 근접 촬영", "시설물마다 가까이서 45도 각도로 촬영"),
    ("상태 조합 전수 커버", "정상·화재·파손·화재+파손 × 6종 조합 모두 촬영"),
    ("좋은 각도부터 촬영", "제한된 시간 내 정보량 많은 각도를 우선 계획"),
]
iy2 = dc_y + Inches(0.44)
for i, (t, d) in enumerate(facility_items):
    add_icon(s6, MSO_SHAPE.OVAL, dc_rx, iy2, Inches(0.26), Inches(0.26), fill=ramp(i, 4))
    add_text(s6, dc_rx, iy2, Inches(0.26), Inches(0.26), str(i + 1), size=9.5, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s6, dc_rx + Inches(0.38), iy2 - Inches(0.02), dc_col_w - Inches(0.38), Inches(0.22),
              t, size=10.3, color=NAVY, bold=True)
    add_text(s6, dc_rx + Inches(0.38), iy2 + Inches(0.21), dc_col_w - Inches(0.38), Inches(0.22),
              d, size=8.8, color=TEXT_GRAY)
    iy2 += dc_step

dc_case_y = Inches(6.30)
dc_case_h = Inches(0.58)
add_rect(s6, MX, dc_case_y, CW, dc_case_h, fill=TEAL_TINT, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
add_pill(s6, MX + Inches(0.2), dc_case_y + Inches(0.15), Inches(1.5), Inches(0.28),
         "핵심 원칙", fill=TEAL, text_color=WHITE, size=8.6)
add_text(s6, MX + Inches(1.9), dc_case_y + Inches(0.13), CW - Inches(2.1), Inches(0.34),
          "한 zone(칸)에는 객체를 최대 1개만 배치해 촬영 — 라벨 경계가 겹치지 않아 바운딩박스 품질이 올라갑니다.",
          size=9.8, color=TEAL_DK, bold=True, anchor=MSO_ANCHOR.MIDDLE)


# =====================================================================
# SLIDE 07 : 핵심 알고리즘① 좌표 변환 파이프라인
# =====================================================================
s7 = prs.slides.add_slide(blank)
set_bg(s7, BG)
header_footer(s7, "07", "src/calibration.py · config/field_config.py 기준", "ALGORITHM · COORDINATE")
title_block(
    s7,
    "ArUco 코너마커 4개 → 호모그래피 → 탑뷰 워핑, 좌표를 한 번에 통일",
    "드론이 프레임마다 흔들려도 매 프레임 재보정 — 이후 모든 crop은 역투영 없이 단순 슬라이싱"
)

steps06 = [
    ("01", "ArUco 마커 검출", "DICT_4X4_50 · 경기장 4모서리 마커",
     "cv2.aruco.ArucoDetector", MSO_SHAPE.DONUT),
    ("02", "서브픽셀 정밀화", "코너 좌표를 소수점 단위로 보정\n(호모그래피 오차 감소 핵심)",
     "cv2.cornerSubPix", MSO_SHAPE.DIAMOND),
    ("03", "호모그래피 계산", "픽셀 → 실좌표(cm) 변환행렬 추정\nRANSAC으로 이상치에 강건하게",
     "cv2.findHomography", MSO_SHAPE.PENTAGON),
    ("04", "탑뷰 워핑", "스케일행렬 × H로 이미지 자체를\n실좌표계와 픽셀 1:1로 정렬",
     "cv2.warpPerspective", MSO_SHAPE.HEXAGON),
]
a_gap = Inches(0.24)
a_w = int((CW - a_gap * 3) / 4)
a_y = Inches(2.55)
a_h = Inches(2.35)
for i, (num, title, desc, detail, shape) in enumerate(steps06):
    x = MX + i * (a_w + a_gap)
    add_rect(s7, x, a_y, a_w, a_h, fill=CARD_BG, line=CARD_BD, line_w=0.75,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    add_text(s7, x + Inches(0.16), a_y + Inches(0.12), Inches(0.6), Inches(0.3),
              num, size=11, color=TEXT_MUTE, bold=True)
    dia = Inches(0.58)
    add_icon(s7, MSO_SHAPE.OVAL, x + (a_w - dia) / 2, a_y + Inches(0.42), dia, dia, fill=ramp(i, 4))
    add_icon(s7, shape, x + (a_w - dia) / 2 + Inches(0.13), a_y + Inches(0.55), dia - Inches(0.26), dia - Inches(0.26),
              fill=WHITE)
    add_text(s7, x + Inches(0.14), a_y + Inches(1.16), a_w - Inches(0.28), Inches(0.28),
              title, size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s7, x + Inches(0.14), a_y + Inches(1.48), a_w - Inches(0.28), Inches(0.55),
              desc, size=8.4, color=TEXT_GRAY, align=PP_ALIGN.CENTER, line_spacing=1.15)
    add_rect(s7, x + Inches(0.12), a_y + a_h - Inches(0.32), a_w - Inches(0.24), Inches(0.24),
              fill=TEAL_TINT, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    add_text(s7, x + Inches(0.12), a_y + a_h - Inches(0.30), a_w - Inches(0.24), Inches(0.2),
              detail, size=7.6, color=TEAL_DK, align=PP_ALIGN.CENTER)

case_y = Inches(5.12)
case_h = Inches(1.42)
add_rect(s7, MX, case_y, CW, case_h, fill=CARD_BG, line=CARD_BD, line_w=0.75,
          shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
add_rect(s7, MX, case_y, Inches(0.07), case_h, fill=ramp(4, 5))
add_pill(s7, MX + Inches(0.24), case_y + Inches(0.16), Inches(1.7), Inches(0.28),
         "실측 보정 사례", fill=ORANGE_BG, text_color=ORANGE_TX, size=8.6)
add_text(s7, MX + Inches(0.24), case_y + Inches(0.55), CW - Inches(0.5), Inches(0.78),
          "사전 가정과 다르게, 실측 결과 ArUco 마커 ID가 0~3이 아닌 1~4였고 코너 대응도 예상과 어긋났습니다. "
          "마커 근접 크롭으로 라벨(\"FA-01\", \"FA-06\" 등)을 직접 확인해 소거법으로 코너를 확정하고, "
          "ARUCO_MARKER_CORNER_INDEX를 재정의해 대응했습니다.",
          size=9.6, color=NAVY, line_spacing=1.3)


# =====================================================================
# SLIDE 08 : 핵심 알고리즘② 탐지·분류 & 프레임 앙상블
# =====================================================================
s8 = prs.slides.add_slide(blank)
set_bg(s8, BG)
header_footer(s8, "08", "src/detection.py · geo_dedup.py 기준", "ALGORITHM · DETECTION")
title_block(
    s8,
    "형태 기술자 점수화 + Union-Find 중복제거로 오탐을 줄입니다",
    "고전 CV 백엔드의 이중 스코어링 판정과, 다중 프레임에 걸친 2단계 앙상블 로직"
)

lt2_x = MX
lt2_w = Inches(5.60)
lt2_y = Inches(2.55)
add_text(s8, lt2_x, lt2_y, lt2_w, Inches(0.3), "classify_blob() 이중 스코어링", size=13, color=NAVY, bold=True)

formula_y = lt2_y + Inches(0.44)
formulas = [
    ("폭파구", "score = 지름차이 × 0.7 + 종횡비차이 × 0.3", "둥글수록 크기가 더 결정적인 단서"),
    ("불발탄", "score = 장축차이 × 0.4 + 종횡비차이 × 0.6", "길쭉한 형태가 더 결정적인 단서"),
]
fy = formula_y
for label, formula, note in formulas:
    add_rect(s8, lt2_x, fy, lt2_w, Inches(0.72), fill=NAVY_DEEP,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    add_pill(s8, lt2_x + Inches(0.16), fy + Inches(0.12), Inches(0.9), Inches(0.24),
             label, fill=TEAL, text_color=WHITE, size=8.2)
    add_text(s8, lt2_x + Inches(0.16), fy + Inches(0.40), lt2_w - Inches(0.32), Inches(0.26),
              formula, size=9.3, color=WHITE, bold=True, font="Consolas")
    fy += Inches(0.84)
add_text(s8, lt2_x, fy + Inches(0.02), lt2_w, Inches(0.5),
          "→ 두 후보군을 동시에 비교해 전역 최소 점수를 선택하므로, 폭파구·불발탄이 서로 중복 집계되지 않습니다.",
          size=9.0, color=TEXT_GRAY, line_spacing=1.25)
add_pill(s8, lt2_x, fy + Inches(0.62), Inches(4.7), Inches(0.32),
         "630개 실측 라벨로 크기 경계값 재보정: 13.1cm / 18.4cm", fill=TEAL_BG2, text_color=TEAL_TX2, size=8.4)

rt2_x = Inches(6.55)
rt2_w = Inches(6.23)
add_text(s8, rt2_x, lt2_y, rt2_w, Inches(0.3), "geo_dedup 2단계 앙상블", size=13, color=NAVY, bold=True)
ensemble_steps = [
    ("실좌표 거리 기반 병합 (Union-Find)",
     "5cm 이내 탐지는 같은 물체로 판정 · 최고 신뢰도를 대표값으로, 좌표는 클러스터 평균으로 스무딩"),
    ("구역별 클래스 재선정",
     "zone(구간) 내 세부 클래스별 신뢰도를 누적 합산 → 가장 높은 점수의 클래스로 최종 확정"),
    ("시계열 다수결 집계",
     "여러 프레임의 (status, confidence) 결과를 Counter 다수결로 집계 (aggregate_temporal_status)"),
]
epy = lt2_y + Inches(0.46)
for i, (h, d) in enumerate(ensemble_steps):
    add_icon(s8, MSO_SHAPE.OVAL, rt2_x, epy, Inches(0.28), Inches(0.28), fill=ramp(i, 3))
    add_text(s8, rt2_x, epy, Inches(0.28), Inches(0.28), str(i + 1), size=10.5, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s8, rt2_x + Inches(0.42), epy - Inches(0.02), rt2_w - Inches(0.42), Inches(0.26),
              h, size=10.5, color=NAVY, bold=True)
    add_text(s8, rt2_x + Inches(0.42), epy + Inches(0.26), rt2_w - Inches(0.42), Inches(0.55),
              d, size=8.6, color=TEXT_GRAY, line_spacing=1.2)
    epy += Inches(0.92)


# =====================================================================
# SLIDE 09 : 핵심 알고리즘③ 활주로 가용길이 & 시설물 강제매핑
# =====================================================================
s9 = prs.slides.add_slide(blank)
set_bg(s9, BG)
header_footer(s9, "09", "src/runway_analysis.py · facility_analysis.py 기준", "ALGORITHM · RUNWAY")
title_block(
    s9,
    "막힌 구간 제외 후 최장 연속 구간 탐색 — O(n) 런렝스 알고리즘",
    "runway_analysis.py의 확정 로직과, 시설물 6슬롯을 항상 채우는 facility_analysis.py 강제 매핑"
)

# 좌측: 활주로 가용길이 시각화
lw_x = MX
lw_w = Inches(5.9)
lw_y = Inches(2.55)
add_text(s9, lw_x, lw_y, lw_w, Inches(0.28), "활주로 가용길이 산출 (예시)", size=12.5, color=NAVY, bold=True)

rw_y = lw_y + Inches(0.5)
rw_h = Inches(0.5)
rw_gap = Inches(0.04)
rw_w = (lw_w - rw_gap * 9) / 10
blocked_idx = {2, 6}  # RW-03, RW-07 (0-indexed)
for i in range(10):
    x = lw_x + i * (rw_w + rw_gap)
    fill = ORANGE_TX if i in blocked_idx else TEAL_BG2
    txc = WHITE if i in blocked_idx else TEAL_TX2
    add_rect(s9, x, rw_y, rw_w, rw_h, fill=fill, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
    add_text(s9, x, rw_y, rw_w, rw_h, f"{i+1:02d}", size=8.5, color=txc, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s9, lw_x, rw_y + rw_h + Inches(0.08), lw_w, Inches(0.22),
          "RW-01~10 (주황 = 폭파구 발견 구간, 청록 = 가용 구간)", size=8.2, color=TEXT_MUTE)

run_y = rw_y + rw_h + Inches(0.42)
add_text(s9, lw_x, run_y, lw_w, Inches(0.24),
          "런(run) 탐색 결과", size=9.5, color=NAVY, bold=True)
run_text = ("[01,02] 2칸(600m)  ·  [04,05,06] 3칸(900m)  ·  [08,09,10] 3칸(900m)\n"
            "→ 동률 시 먼저 나오는 런 채택: 최장 가용길이 = 900m")
add_text(s9, lw_x, run_y + Inches(0.26), lw_w, Inches(0.55),
          run_text, size=9.0, color=TEXT_GRAY, line_spacing=1.3)

algo_note_y = run_y + Inches(0.92)
add_rect(s9, lw_x, algo_note_y, lw_w, Inches(0.85), fill=TEAL_TINT,
          shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
add_text(s9, lw_x + Inches(0.18), algo_note_y + Inches(0.12), lw_w - Inches(0.36), Inches(0.62),
          "알고리즘: 구간을 순서대로 훑으며 막힌 구간에서 현재 런을 flush,\n"
          "가용 구간이면 런에 누적 — 1회 순회(O(n))로 최장 연속 구간을 확정",
          size=8.6, color=TEAL_DK, line_spacing=1.3)

# 우측: 시설물 6슬롯 강제매핑
rw2_x = Inches(6.85)
rw2_w = Inches(5.93)
add_text(s9, rw2_x, lw_y, rw2_w, Inches(0.28), "시설물 6슬롯 강제 매핑", size=12.5, color=NAVY, bold=True)

fac_status = [("FA-01", "normal", TEAL_BG2, TEAL_TX2), ("FA-02", "fire", ORANGE_BG, ORANGE_TX),
              ("FA-03", "unconfirmed", CARD_BD, TEXT_MUTE), ("FA-04", "normal", TEAL_BG2, TEAL_TX2),
              ("FA-05", "destroy", ORANGE_BG, ORANGE_TX), ("FA-06", "normal", TEAL_BG2, TEAL_TX2)]
fg_y = lw_y + Inches(0.5)
fg_gap = Inches(0.14)
fg_w = (rw2_w - fg_gap * 2) / 3
fg_h = Inches(0.85)
for i, (slot, status, fill, txc) in enumerate(fac_status):
    col = i % 3
    row = i // 3
    x = rw2_x + col * (fg_w + fg_gap)
    y = fg_y + row * (fg_h + fg_gap)
    add_rect(s9, x, y, fg_w, fg_h, fill=CARD_BG, line=CARD_BD, line_w=0.75,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    add_rect(s9, x, y, fg_w, Inches(0.06), fill=fill)
    add_text(s9, x + Inches(0.12), y + Inches(0.14), fg_w - Inches(0.24), Inches(0.24),
              slot, size=10.5, color=NAVY, bold=True)
    add_pill(s9, x + Inches(0.12), y + Inches(0.44), fg_w - Inches(0.24), Inches(0.26),
             status, fill=fill, text_color=txc, size=8.0)

fac_note_y = fg_y + fg_h * 2 + fg_gap + Inches(0.16)
add_text(s9, rw2_x, fac_note_y, rw2_w, Inches(0.8),
          "탐지 성공 → status 그대로 기록 / 탐지 실패 → 슬롯을 비우지 않고 \"unconfirmed\"로 명시.\n"
          "위치가 고정되어 있다는 사전 지식을 활용해 항상 6개 슬롯을 채워, \"시설물 누락\" 감점을 원천 차단합니다.",
          size=9.0, color=TEXT_GRAY, line_spacing=1.3)


# =====================================================================
# SLIDE 10 : 데이터 신뢰성 & 자동 검증
# =====================================================================
s10 = prs.slides.add_slide(blank)
set_bg(s10, BG)
header_footer(s10, "10", "src/validator.py · report_generator.py · requirements.txt 기준", "DATA RELIABILITY")
title_block(
    s10,
    "전송 전 6가지 자동 QA로 사람의 실수를 코드가 대신 잡습니다",
    "validator.py 6종 체크 + LLM 안전장치(환각 차단) + 검증된 오픈소스 기술 스택"
)

checks = [
    "mission_code 일관성 — 8개 JSON 파일 전부 동일해야 함",
    "crater_count ≤ crater_detect 총 개수 (초과 불가)",
    "crater_detect / uxo_detect 건수가 서버 허용 범위 안에 있는지",
    "시설물 6슬롯 전부 존재 — 누락 감점 차단 핵심 체크",
    "uxo_count ≤ uxo_detect 총 개수 (초과 불가)",
    "활주로 가용길이 물리적 범위 + report 글자수 50~100자 준수",
]
chk_y = Inches(2.55)
chk_gap = Inches(0.18)
chk_w = (CW - chk_gap) / 2
chk_h = Inches(0.62)
for i, text in enumerate(checks):
    col = i % 2
    row = i // 2
    x = MX + col * (chk_w + chk_gap)
    y = chk_y + row * (chk_h + Inches(0.12))
    add_rect(s10, x, y, chk_w, chk_h, fill=CARD_BG, line=CARD_BD, line_w=0.75,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    add_icon(s10, MSO_SHAPE.OVAL, x + Inches(0.14), y + chk_h / 2 - Inches(0.13), Inches(0.26), Inches(0.26),
              fill=ramp(i, 6))
    add_text(s10, x + Inches(0.14), y + chk_h / 2 - Inches(0.13), Inches(0.26), Inches(0.26),
              str(i + 1), size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s10, x + Inches(0.52), y, chk_w - Inches(0.66), chk_h,
              text, size=9.6, color=NAVY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)

llm_y = chk_y + 3 * (chk_h + Inches(0.12)) + Inches(0.06)
add_rect(s10, MX, llm_y, CW, Inches(0.72), fill=NAVY_DEEP, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
add_pill(s10, MX + Inches(0.2), llm_y + Inches(0.13), Inches(1.7), Inches(0.26),
         "LLM 안전장치", fill=TEAL, text_color=WHITE, size=8.4)
add_text(s10, MX + Inches(2.05), llm_y + Inches(0.10), CW - Inches(2.3), Inches(0.52),
          "보고시각·개수·가용길이·상태는 코드가 먼저 확정해 \"사실\"로 주입 — LLM은 문장 생성만 담당(수치 환각 차단). "
          "응답 타임아웃 8초, 실패 시 결정론적 템플릿으로 자동 폴백.",
          size=9.1, color=WHITE, line_spacing=1.2, anchor=MSO_ANCHOR.MIDDLE)

stack_y = llm_y + Inches(0.72) + Inches(0.24)
add_text(s10, MX, stack_y, CW, Inches(0.26), "검증된 기술 스택", size=11, color=NAVY, bold=True)
techs = ["Python 3.12", "OpenCV (cv2.aruco)", "NumPy", "Ultralytics YOLO11n",
         "PyTorch + CUDA", "watchdog", "Ollama qwen2.5:7b", "requests"]
tx = MX
ty = stack_y + Inches(0.36)
for i, t in enumerate(techs):
    w = Inches(0.22 + 0.105 * len(t))
    if tx + w > MX + CW:
        tx = MX
        ty += Inches(0.4)
    add_pill(s10, tx, ty, w, Inches(0.32), t, fill=TEAL_TINT, text_color=TEAL_DK, size=8.6)
    tx += w + Inches(0.14)


# =====================================================================
# SLIDE 11 : 임무 수행 전략  (발표평가 "임무 수행 전략" 대응)
# =====================================================================
s11 = prs.slides.add_slide(blank)
set_bg(s11, BG)
header_footer(s11, "11", "공통운영규칙 경기시간 규정(대기1분·준비1분·임무3분) 기준", "MISSION STRATEGY")
title_block(
    s11,
    "대기 1분 · 준비 1분 · 임무 3분 — 초 단위로 설계한 실행 전략",
    "경기 규정(총 5분·300초)에 맞춰 무엇을 언제 하는지 미리 확정했습니다"
)

bar_y = Inches(2.62)
bar_h = Inches(0.6)
seg_gap = Inches(0.12)
total_units = 1 + 1 + 3
unit_w = (CW - seg_gap * 2) / total_units
segments = [
    ("대기 60초", 1, TEAL_TINT, TEAL_DK,
     "지정 위치 이동 · 노트북/드론 배치\n프로그램 실행·조작 금지 구간"),
    ("준비 60초", 1, ramp(2, 5), WHITE,
     "드론 조종앱 · AI 프로그램 실행\n종료 전 이륙 완료 필수(자동실격 주의)"),
    ("임무수행 180초", 3, NAVY_DEEP, WHITE,
     "이륙~착륙까지 watchdog 자동감지\n→배치추론→8종 JSON 순차 전송"),
]
sx = MX
for label, units, fill, txc, desc in segments:
    w = unit_w * units
    add_rect(s11, sx, bar_y, w, bar_h, fill=fill, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    add_text(s11, sx, bar_y, w, bar_h, label, size=13, color=txc, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s11, sx, bar_y + bar_h + Inches(0.14), w, Inches(0.6),
              desc, size=9.1, color=TEXT_GRAY, align=PP_ALIGN.CENTER, line_spacing=1.15)
    sx += w + seg_gap

stat_y = Inches(4.75)
stat_h = Inches(1.85)
stats6 = [
    ("배치추론 1회 호출", "20개 zone·6개 ROI를 각각 한 번의\nmodel.predict(list)로 처리"),
    ("watchdog 자동 파이프라인", "영상 저장 즉시 프레임 추출→추론 시작\n수동 선택 없이 전 과정 자동화"),
    ("순서 무관 순차 전송", "8종 JSON, 준비되는 대로 즉시 전송\n중복 전송 허용 규정을 그대로 활용"),
    ("이륙 후 무조작 원칙", "조작·재실행 금지 규정 위반 원천 차단\n점검은 반드시 준비시간 내 완료"),
]
sn = len(stats6)
sgap = Inches(0.24)
sw = int((CW - sgap * (sn - 1)) / sn)
add_text(s11, MX, stat_y - Inches(0.38), CW, Inches(0.3), "3분 대응 핵심 전략", size=13, color=NAVY, bold=True)
for i, (head, desc) in enumerate(stats6):
    x = MX + i * (sw + sgap)
    add_rect(s11, x, stat_y, sw, stat_h, fill=CARD_BG, line=CARD_BD, line_w=0.75,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07)
    add_rect(s11, x, stat_y, Inches(0.06), stat_h, fill=ramp(i, 4))
    add_text(s11, x + Inches(0.2), stat_y + Inches(0.16), sw - Inches(0.36), Inches(0.5),
              head, size=11, color=NAVY, bold=True, line_spacing=1.05)
    add_text(s11, x + Inches(0.2), stat_y + Inches(0.68), sw - Inches(0.36), Inches(1.0),
              desc, size=9.0, color=TEXT_GRAY, line_spacing=1.2)


# =====================================================================
# SLIDE 12 : 마무리 (요약 & 질의응답 대비)
# =====================================================================
s12 = prs.slides.add_slide(blank)
set_bg(s12, BG)
header_footer(s12, "12", "readme.md 알려진 한계 섹션 기준", "SUMMARY")
title_block(
    s12,
    "정리하면",
    "27시간의 시행착오로 완성한 3가지 핵심 — 좌표 통일 · 무중단 설계 · 규정 정면 대응"
)

sum_y = Inches(2.55)
sum_h = Inches(1.75)
summary_cards = [
    ("좌표계 통일", MSO_SHAPE.DIAMOND, "ArUco 호모그래피로 모든 탐지·분류의\n기준 좌표를 하나로 정렬"),
    ("무중단 설계", MSO_SHAPE.HEXAGON, "시설물 6슬롯 강제 매핑 + 로컬LLM·\n오프라인 폴백으로 임무 절대 중단 없음"),
    ("규정 정면 대응", MSO_SHAPE.UP_ARROW, "3분 제한 · 100자 제약 · 8종 JSON 스펙\n모두 코드 로직으로 직접 구현"),
]
sn2 = len(summary_cards)
sgap2 = Inches(0.24)
sw2 = int((CW - sgap2 * (sn2 - 1)) / sn2)
for i, (head, shape, desc) in enumerate(summary_cards):
    x = MX + i * (sw2 + sgap2)
    add_rect(s12, x, sum_y, sw2, sum_h, fill=CARD_BG, line=CARD_BD, line_w=0.75,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07)
    dia = Inches(0.5)
    add_icon(s12, MSO_SHAPE.OVAL, x + Inches(0.18), sum_y + Inches(0.18), dia, dia, fill=ramp(i, 3))
    add_icon(s12, shape, x + Inches(0.30), sum_y + Inches(0.30), dia - Inches(0.24), dia - Inches(0.24), fill=WHITE)
    add_text(s12, x + Inches(0.18), sum_y + Inches(0.78), sw2 - Inches(0.36), Inches(0.3),
              head, size=13, color=NAVY, bold=True)
    add_text(s12, x + Inches(0.18), sum_y + Inches(1.10), sw2 - Inches(0.36), Inches(0.6),
              desc, size=9.1, color=TEXT_GRAY, line_spacing=1.2)

# 알려진 한계 & 향후 계획
lim_y = Inches(4.55)
lim_h = Inches(0.95)
add_rect(s12, MX, lim_y, CW, lim_h, fill=TEAL_TINT, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
add_text(s12, MX + Inches(0.22), lim_y + Inches(0.13), Inches(3.4), Inches(0.28),
          "알려진 한계 & 향후 계획", size=11, color=TEAL_DK, bold=True)
lim_text = ("· zone 경계에 걸친 물체는 그리드 분할 특성상 분리 탐지될 수 있음 → overlap(SAHI) 적용 검토\n"
            "· 조명 조건에 민감한 화재 판정 임계값 → 실전 데이터 축적으로 지속 보정 예정")
add_text(s12, MX + Inches(0.22), lim_y + Inches(0.42), CW - Inches(0.44), Inches(0.5),
          lim_text, size=9.5, color=NAVY, line_spacing=1.25)

# 클로징
add_text(s12, MX, Inches(5.85), CW, Inches(0.7), "감사합니다", size=32, color=NAVY, bold=True,
          align=PP_ALIGN.CENTER)
add_text(s12, MX, Inches(6.55), CW, Inches(0.35), "Team bongbong7 · 제8회 공군 해커톤 AI경진대회",
          size=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bongbong7_본선발표자료.pptx")
prs.save(out_path)
print("SAVED:", out_path)
